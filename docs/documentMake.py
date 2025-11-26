from pptx import Presentation
from pptx.util import Inches
import os

imgFolder = "./docs/img"

ppt = Presentation()
#ppt.slide_height = Inches(7.5)
#ppt.slide_width =  Inches(13.33)

files = os.listdir(imgFolder)
for idx,file in enumerate(files):
    FName, ext = os.path.splitext(file)
    if ext == ".png":
        files[idx] = file.strip(".png")


files = map(int, files)
files = sorted(files)


for idx,file in enumerate(files):
    files[idx] = f"{file}.png"



for img in files:
    if img.lower().endswith(".png"):
        slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        print(img)
        
        slide.shapes.add_picture(os.path.join(imgFolder, img),
                                  left=0, 
                                  top=0, 
                                  height=ppt.slide_height,
                                  width=ppt.slide_width
                                )

ppt.save("./docs/Documentation.pptx")